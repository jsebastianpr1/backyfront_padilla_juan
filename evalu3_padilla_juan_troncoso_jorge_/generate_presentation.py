import sys
import subprocess

# Auto-install python-pptx if it is not present
try:
    import pptx
except ImportError:
    print("Installing python-pptx library...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Style variables
bg_color = RGBColor(10, 15, 29)       # #0A0F1D (Very dark blue background)
card_color = RGBColor(22, 27, 48)     # #161B30 (Card Background)
accent_cyan = RGBColor(0, 240, 255)   # #00F0FF (Electric Cyan)
accent_pink = RGBColor(255, 0, 127)   # #FF007F (Neon Magenta)
text_white = RGBColor(255, 255, 255)  # White
text_grey = RGBColor(160, 174, 192)   # Light grey
code_bg_color = RGBColor(14, 18, 34)  # Code block background
code_text_color = RGBColor(56, 189, 248) # Code text (Sky Blue)

# Helper function to create blank slide with dark background
def create_slide(title_text, category="PROVIEMPLEA - BACKEND"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = bg_color
    bg.line.fill.background()
    
    # Category Text
    cat_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(10), Inches(0.4))
    tf_cat = cat_box.text_frame
    tf_cat.word_wrap = True
    p_cat = tf_cat.paragraphs[0]
    p_cat.text = category.upper()
    p_cat.font.name = "Arial"
    p_cat.font.size = Pt(11)
    p_cat.font.bold = True
    p_cat.font.color.rgb = accent_cyan
    
    # Title Text
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.7), Inches(12), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = "Trebuchet MS"
    p_title.font.size = Pt(26)
    p_title.font.bold = True
    p_title.font.color.rgb = text_white
    
    # Accent Line under Title
    accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.5), Inches(12.133), Inches(0.04))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = accent_pink
    accent_line.line.fill.background()
    
    return slide

# Helper to add content card
def add_card(slide, left, top, width, height, title_text="", border_color=None):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = card_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
    else:
        card.line.fill.background()
        
    if title_text:
        title_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Arial"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = accent_cyan
        
    return card

# Helper to add code block
def add_code_block(slide, left, top, width, height, code_text):
    code_card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    code_card.fill.solid()
    code_card.fill.fore_color.rgb = code_bg_color
    code_card.line.color.rgb = RGBColor(0x3B, 0x42, 0x52)
    code_card.line.width = Pt(1)
    
    tb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1), width - Inches(0.3), height - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    
    lines = code_text.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = "Consolas"
        p.font.size = Pt(10.5)
        p.font.color.rgb = code_text_color

# Helper to add bullets inside cards
def add_bullets(slide, left, top, width, height, bullets_list, font_size=13):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        
        if ":" in item and not item.startswith("http"):
            parts = item.split(":", 1)
            p.text = parts[0] + ":"
            p.font.bold = True
            p.font.name = "Arial"
            p.font.size = Pt(font_size)
            p.font.color.rgb = text_white
            
            run = p.add_run()
            run.text = parts[1]
            run.font.bold = False
            run.font.name = "Arial"
            run.font.size = Pt(font_size)
            run.font.color.rgb = text_grey
        else:
            p.text = item
            p.font.name = "Arial"
            p.font.size = Pt(font_size)
            p.font.color.rgb = text_grey

# ==============================================================================
# SLIDE 1: PORTADA
# ==============================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg1.fill.solid()
bg1.fill.fore_color.rgb = bg_color
bg1.line.fill.background()

tb1 = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(2.0))
tf1 = tb1.text_frame
tf1.word_wrap = True
p = tf1.paragraphs[0]
p.text = "PROVIEMPLEA BACKEND"
p.font.name = "Trebuchet MS"
p.font.size = Pt(46)
p.font.bold = True
p.font.color.rgb = accent_cyan

p2 = tf1.add_paragraph()
p2.text = "Guía de Puesta en Marcha: API CRUD y Servidor SQLite"
p2.font.name = "Arial"
p2.font.size = Pt(22)
p2.font.color.rgb = text_white
p2.space_before = Pt(10)

line1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(4.2), Inches(6.0), Inches(0.08))
line1.fill.solid()
line1.fill.fore_color.rgb = accent_pink
line1.line.fill.background()

meta_box = slide1.shapes.add_textbox(Inches(1.0), Inches(5.2), Inches(8.0), Inches(1.5))
tf_meta = meta_box.text_frame
tf_meta.word_wrap = True

p_m1 = tf_meta.paragraphs[0]
p_m1.text = "Desarrollado para la Evaluación 3 de Integración Frontend-Backend"
p_m1.font.name = "Arial"
p_m1.font.size = Pt(13)
p_m1.font.color.rgb = text_grey

p_m2 = tf_meta.add_paragraph()
p_m2.text = "Autores: Juan Padilla & Jorge Troncoso"
p_m2.font.bold = True
p_m2.font.name = "Arial"
p_m2.font.size = Pt(14)
p_m2.font.color.rgb = accent_cyan
p_m2.space_before = Pt(5)

p_m3 = tf_meta.add_paragraph()
p_m3.text = "Municipalidad de Providencia • Framework Laravel 12 + SQLite"
p_m3.font.name = "Arial"
p_m3.font.size = Pt(12)
p_m3.font.color.rgb = text_grey
p_m3.space_before = Pt(5)

# ==============================================================================
# SLIDE 2: EL ECOSISTEMA
# ==============================================================================
slide2 = create_slide("El Ecosistema Backend: Laravel + SQLite", "01. CONTEXTO Y ARQUITECTURA")

add_card(slide2, Inches(0.6), Inches(1.8), Inches(5.8), Inches(5.0), "La Base de Datos SQLite")
bullets_left2 = [
    "Motor Serverless: SQLite es una base de datos integrada directamente en un archivo local sin requerir un demonio de servidor activo.",
    "Ruta del Archivo: Se almacena localmente en la carpeta del proyecto en `database/database.sqlite`.",
    "Cero Administración: No requiere instalar MySQL, PostgreSQL ni configurar puertos o credenciales complejas.",
    "Portabilidad Extrema: Facilita que todo el entorno se ejecute al instante copiando los archivos del proyecto."
]
add_bullets(slide2, Inches(0.8), Inches(2.5), Inches(5.4), Inches(4.0), bullets_left2)

add_card(slide2, Inches(6.9), Inches(1.8), Inches(5.8), Inches(5.0), "La API CRUD en Laravel")
bullets_right2 = [
    "Framework Moderno: Desarrollado sobre PHP Laravel 12 con soporte integrado para controladores de API RESTful.",
    "Controlador API: Administrado por `UserController.php`, con soporte completo para operaciones CRUD básicas.",
    "Perfil Completo: El método `toBlindArray()` expone todos los datos profesionales y demográficos del candidato.",
    "Capa de Caché: Utiliza base de datos para almacenar temporalmente resultados de consultas pesadas (< 15ms)."
]
add_bullets(slide2, Inches(7.1), Inches(2.5), Inches(5.4), Inches(4.0), bullets_right2)

# ==============================================================================
# SLIDE 3: PASO 1 - INICIALIZAR SQLITE
# ==============================================================================
slide3 = create_slide("Paso 1: Inicialización de SQLite y Entorno", "02. CONFIGURACIÓN DEL ENTORNO")

add_card(slide3, Inches(0.6), Inches(1.8), Inches(6.5), Inches(5.0), "Comandos de Consola (PowerShell)")
bullets_left3 = [
    "1. Configurar Variable de Entorno: Si PHP no está registrado globalmente, agregue temporalmente la ruta del ejecutable PHP 8.3.",
    "2. Ejecutar Migraciones y Semillas: Reconstruye las tablas de base de datos desde cero y carga datos de prueba iniciales."
]
add_bullets(slide3, Inches(0.8), Inches(2.4), Inches(6.1), Inches(1.8), bullets_left3)

code_text_3 = """# Registrar PHP local en la sesión de terminal
$env:PATH += ";C:\\Users\\juans\\AppData\\Local\\Microsoft\\WinGet\\Packages\\PHP.PHP.8.3_Microsoft.Winget.Source_8wekyb3d8bbwe"

# Ejecutar migraciones limpias y sembrar datos
php artisan migrate:fresh --seed"""
add_code_block(slide3, Inches(0.8), Inches(4.3), Inches(6.1), Inches(2.2), code_text_3)

add_card(slide3, Inches(7.4), Inches(1.8), Inches(5.3), Inches(5.0), "¿Qué ocurre tras bambalinas?")
bullets_right3 = [
    "Creación del Archivo: Laravel crea automáticamente el archivo `database.sqlite` (si no existiera).",
    "Esquema de Tablas: Se generan la tabla `users` para candidatos y la tabla `cache` para la gestión de rendimiento.",
    "Población de Datos: El sembrador `DatabaseSeeder.php` genera 5 perfiles diversos con competencias técnicas y casts JSON.",
    "Consistencia de Datos: Garantiza que la base de datos esté lista para recibir solicitudes CRUD del frontend."
]
add_bullets(slide3, Inches(7.6), Inches(2.4), Inches(4.9), Inches(4.1), bullets_right3)

# ==============================================================================
# SLIDE 4: PASO 2 - LEVANTAR EL SERVIDOR DE LA API
# ==============================================================================
slide4 = create_slide("Paso 2: Levantar el Servidor de la API CRUD", "03. ENCENDIDO DE LA API")

add_card(slide4, Inches(0.6), Inches(1.8), Inches(6.5), Inches(5.0), "Ejecución del Servidor HTTP")
bullets_left4 = [
    "Servidor de Desarrollo: Iniciamos el servidor web local integrado de Laravel (`Artisan Serve`).",
    "Puerto Especificado: Se define el puerto 8080 para evitar conflictos con otros servicios locales.",
    "Conexión con SQLite: Cada petición HTTP entrante activará Eloquent para consultar de forma transparente el archivo SQLite local."
]
add_bullets(slide4, Inches(0.8), Inches(2.4), Inches(6.1), Inches(2.0), bullets_left4)

code_text_4 = """# Ejecutar en la carpeta raíz del proyecto
php artisan serve --port=8080"""
add_code_block(slide4, Inches(0.8), Inches(4.5), Inches(6.1), Inches(1.8), code_text_4)

add_card(slide4, Inches(7.4), Inches(1.8), Inches(5.3), Inches(5.0), "Información del Servicio")
bullets_right4 = [
    "Punto de Acceso Local: La API CRUD estará disponible para recibir llamadas en `http://localhost:8080/api/usuarios`",
    "SQLite en Acción: No es necesario ejecutar ningún comando como 'start sqlite' o similar. SQLite arranca automáticamente al consultar el archivo.",
    "Validación Rápida: Al abrir la URL anterior en un navegador, debería ver una respuesta JSON con la lista de candidatos con su perfil completo."
]
add_bullets(slide4, Inches(7.6), Inches(2.4), Inches(4.9), Inches(4.1), bullets_right4)

# ==============================================================================
# SLIDE 5: ENDPOINTS DISPONIBLES
# ==============================================================================
slide5 = create_slide("Endpoints CRUD Disponibles de la API", "04. CAPA DE SERVICIO")

add_card(slide5, Inches(0.6), Inches(1.8), Inches(12.133), Inches(5.0), "Catálogo de Rutas de la Vitrina de Empleo")
bullets_full5 = [
    "GET /api/usuarios : Obtiene los candidatos activos. Retorna perfiles de candidatos completos (incluyendo nombres, correos y comunas) para una visualización transparente.",
    "GET /api/usuarios?role=admin : Ruta administrativa opcional para que el personal autorizado audite o administre los perfiles completos de los candidatos.",
    "POST /api/usuarios : Permite registrar nuevos candidatos en el sistema. Valida obligatoriamente campos sensibles y formaciones profesionales.",
    "PATCH /api/usuarios/{id}/status : Enciende o apaga la visibilidad del perfil (is_active). Útil cuando un candidato consigue empleo y desea suspender temporalmente su visualización.",
    "DELETE /api/usuarios/{id} : Remueve permanentemente un postulante del sistema y realiza una invalidación instantánea de la caché."
]
add_bullets(slide5, Inches(0.8), Inches(2.4), Inches(11.7), Inches(4.2), bullets_full5, font_size=12)

# ==============================================================================
# SLIDE 6: CONSUMO Y PRUEBAS
# ==============================================================================
slide6 = create_slide("Consumo y Pruebas en PowerShell", "05. PRUEBAS DE INTEGRACIÓN")

add_card(slide6, Inches(0.6), Inches(1.8), Inches(5.8), Inches(5.0), "1. Listar Vitrina de Empleo (GET)")
bullets_l6 = [
    "Consulta Inicial: Solicita la lista de candidatos activos. Retorna datos profesionales y sociodemográficos completos.",
    "Prueba de Endpoint: Ejecute este comando en una terminal de PowerShell secundaria:"
]
add_bullets(slide6, Inches(0.8), Inches(2.4), Inches(5.4), Inches(1.8), bullets_l6)

code_get = """# Consultar candidatos activos (Perfil Completo)
Invoke-RestMethod -Uri "http://localhost:8080/api/usuarios" `
  -Method GET"""
add_code_block(slide6, Inches(0.8), Inches(4.3), Inches(5.4), Inches(2.0), code_get)

add_card(slide6, Inches(6.9), Inches(1.8), Inches(5.8), Inches(5.0), "2. Apagar Visibilidad (PATCH)")
bullets_r6 = [
    "Toggle de Estado: Envía una actualización parcial indicando 'is_active = false' para suspender a un candidato.",
    "Comando de Consola: Ejecute lo siguiente en PowerShell:"
]
add_bullets(slide6, Inches(7.1), Inches(2.4), Inches(5.4), Inches(1.8), bullets_r6)

code_patch = """# Definir cuerpo JSON y apagar candidato ID 4
$body = @{ is_active = $false } | ConvertTo-Json

Invoke-RestMethod -Uri "http://localhost:8080/api/usuarios/4/status" `
  -Method PATCH -Body $body -ContentType "application/json\""""
add_code_block(slide6, Inches(7.1), Inches(4.3), Inches(5.4), Inches(2.0), code_patch)

# ==============================================================================
# SLIDE 7: RENDIMIENTO Y SWAGGER
# ==============================================================================
slide7 = create_slide("Rendimiento Extremo y Documentación", "06. OPTIMIZACIÓN Y DOCUMENTACIÓN")

add_card(slide7, Inches(0.6), Inches(1.8), Inches(5.8), Inches(5.0), "Estrategia Cache-Aside")
bullets_l7 = [
    "Almacenamiento Local: Respuestas de búsquedas y filtros se indexan con una clave única (`users:all:{hash_de_filtros}`).",
    "Cache Hit: Las búsquedas en caché tardan menos de 15 milisegundos, liberando de carga al servidor SQLite.",
    "Invalidación Instantánea: Cualquier cambio (POST, PUT, DELETE, PATCH) invoca `Cache::flush()`, asegurando que las empresas vean cambios en tiempo real."
]
add_bullets(slide7, Inches(0.8), Inches(2.4), Inches(5.4), Inches(4.0), bullets_l7)

add_card(slide7, Inches(6.9), Inches(1.8), Inches(5.8), Inches(5.0), "Documentación Swagger")
bullets_r7 = [
    "Archivo Local: El archivo `swagger.yaml` describe toda la API en español con ejemplos estructurados de peticiones.",
    "Visualización Interactiva: Puede copiar su contenido y pegarlo en `editor.swagger.io` para interactuar en vivo.",
    "Simulador Integrado: Permite probar parámetros de filtrado (habilidades, certificaciones), probar límites de peticiones y verificar respuestas de error."
]
add_bullets(slide7, Inches(7.1), Inches(2.4), Inches(5.4), Inches(4.0), bullets_r7)

# ==============================================================================
# SLIDE 8: CONCLUSIÓN / CIERRE
# ==============================================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
bg8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg8.fill.solid()
bg8.fill.fore_color.rgb = bg_color
bg8.line.fill.background()

tb8 = slide8.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(2.0))
tf8 = tb8.text_frame
tf8.word_wrap = True
p = tf8.paragraphs[0]
p.text = "¡TODO LISTO!"
p.font.name = "Trebuchet MS"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = accent_cyan

p2 = tf8.add_paragraph()
p2.text = "La API CRUD de ProviEmplea y el Servidor SQLite están listos para ser operados."
p2.font.name = "Arial"
p2.font.size = Pt(20)
p2.font.color.rgb = text_white
p2.space_before = Pt(10)

line8 = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(4.2), Inches(8.0), Inches(0.06))
line8.fill.solid()
line8.fill.fore_color.rgb = accent_pink
line8.line.fill.background()

meta_box8 = slide8.shapes.add_textbox(Inches(1.0), Inches(4.8), Inches(10.0), Inches(2.0))
tf_meta8 = meta_box8.text_frame
tf_meta8.word_wrap = True

p_m8 = tf_meta8.paragraphs[0]
p_m8.text = "Pasos finales recomendados:"
p_m8.font.name = "Arial"
p_m8.font.bold = True
p_m8.font.size = Pt(14)
p_m8.font.color.rgb = accent_cyan

bullets_end = [
    "1. Abra una terminal PowerShell en la raíz del proyecto para inicializar la base de datos.",
    "2. Ejecute el comando `php artisan serve --port=8080` para encender el servidor.",
    "3. Abra `http://localhost:8080/api/usuarios` en su navegador para validar el funcionamiento.",
    "4. ¡Éxito en su presentación de integración frontend-backend!"
]
for item in bullets_end:
    p_b = tf_meta8.add_paragraph()
    p_b.text = item
    p_b.font.name = "Arial"
    p_b.font.size = Pt(13)
    p_b.font.color.rgb = text_grey
    p_b.space_before = Pt(4)

# Save the presentation
output_path = "presentacion_api_sqlite_updated.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
